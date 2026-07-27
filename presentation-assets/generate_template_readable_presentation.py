from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(
    r"C:\Users\Samsung\Downloads\Vibeathon_6.0_Vibecoding_Hackathon_July_2026_Idea_Submission_Template.pptx"
)
OUTPUT = ROOT / "DINE_AI_Vibeathon_Template_Filled_READABLE.pptx"

PURPLE = RGBColor(52, 29, 66)
INK = RGBColor(20, 21, 34)
MUTED = RGBColor(78, 83, 96)
WHITE = RGBColor(255, 255, 255)
CITRUS = RGBColor(216, 250, 73)
CORAL = RGBColor(230, 78, 82)
MINT = RGBColor(34, 166, 139)
LINE = RGBColor(232, 229, 223)


def clear_text(shape):
    if shape.has_text_frame:
        shape.text_frame.clear()


def write_textbox(slide, text, left, top, width, height, size=18, color=INK, bold=False, center=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def panel(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    return shape


def accent_panel(slide, title, text, left, top, width, height, accent=CORAL):
    panel(slide, left, top, width, height)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.09), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent
    write_textbox(slide, title, left + Inches(0.25), top + Inches(0.12), width - Inches(0.35), Inches(0.32), 15, PURPLE, True)
    write_textbox(slide, text, left + Inches(0.25), top + Inches(0.52), width - Inches(0.42), height - Inches(0.58), 11, INK)


def set_template_title(slide, title):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            clear_text(shape)
            write_textbox(slide, title, shape.left, shape.top, shape.width, shape.height, 32, PURPLE, True, True)
            return


def clear_template_content(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            clear_text(shape)


def team_logo(slide):
    panel(slide, Inches(0.32), Inches(0.16), Inches(1.25), Inches(0.55))
    write_textbox(slide, "DINE AI", Inches(0.36), Inches(0.26), Inches(1.15), Inches(0.25), 14, PURPLE, True, True)


def main():
    prs = Presentation(TEMPLATE)

    content = [
        {
            "title": "DINE AI",
            "body": (
                "TITLE PAGE\n\n"
                "Team Name: Byte Crunchers\n"
                "Team Leader Name: Riddhi Bathla\n"
                "College Name: Institute of Management Studies\n"
                "Year & Department: 3rd year & BCA\n\n"
                "Problem Statement / Project Title:\n"
                "DINE AI - Smart Restaurant Operations Platform\n\n"
                "A real-time operating layer for guests, service teams, kitchen, inventory, billing and dietary safety."
            ),
        },
        {
            "title": "CURRENT PROBLEM",
            "body": (
                "Problem Overview:\n"
                "Restaurants still rely on disconnected tools for ordering, waitlists, kitchen tickets, inventory and billing.\n\n"
                "Who Is Affected:\n"
                "Guests, servers, kitchen staff, managers and restaurant owners.\n\n"
                "Pain Points:\n"
                "- Wait times are unclear for guests\n"
                "- Dietary preferences can be missed during handoffs\n"
                "- Inventory shortages are noticed too late\n"
                "- Managers lack one live view of shift health\n\n"
                "Why It Matters:\n"
                "A small communication miss can delay service, waste stock or create a safety risk."
            ),
        },
        {
            "title": "PROPOSED SOLUTION",
            "body": (
                "Solution Overview:\n"
                "DINE AI connects guests, service, kitchen, inventory, billing and analytics in one workspace.\n\n"
                "How It Works:\n"
                "- Guests browse dishes and select preferences\n"
                "- Staff manage tables, waitlist and service tasks\n"
                "- Kitchen sees tickets with timing and safety context\n"
                "- Managers track live signals and stock pressure\n\n"
                "Key Features:\n"
                "- Guest menu\n"
                "- Waitlist\n"
                "- Service floor\n"
                "- Kitchen line\n"
                "- Pantry impact\n"
                "- SafePlate dietary relay\n"
                "- Billing and analytics"
            ),
        },
        {
            "title": "TECHNICAL APPROACH",
            "body": (
                "Technologies Used:\n"
                "- Frontend: Next.js, React, TypeScript, Tailwind CSS\n"
                "- Backend: Node.js app.js\n"
                "- API: Separate routes and controllers\n"
                "- UI: Lucide icons and responsive workspace styling\n\n"
                "System Architecture:\n"
                "frontend/app -> components -> lib/demo-data\n"
                "backend/app.js -> routes -> controllers -> data\n\n"
                "Methodology / Workflow:\n"
                "1. User opens role workspace\n"
                "2. UI displays operational data\n"
                "3. Staff acts on focused cards\n"
                "4. Backend API can serve menu, operations and dashboard data"
            ),
        },
        {
            "title": "USE CASES & IMPACT",
            "body": (
                "Key Use Cases:\n"
                "- Guest ordering with dietary preferences\n"
                "- Staff table service and waitlist management\n"
                "- Kitchen prioritization\n"
                "- Inventory visibility before stockouts\n"
                "- Manager shift analytics\n\n"
                "Target Users:\n"
                "Restaurant managers, servers, kitchen teams, guests and owners.\n\n"
                "Expected Impact:\n"
                "- Faster service handoffs\n"
                "- Better allergen communication\n"
                "- Fewer stock surprises\n"
                "- Clearer manager decisions"
            ),
        },
        {
            "title": "FUTURE SCOPE & CONCLUSION",
            "body": (
                "Future Enhancements:\n"
                "- POS and payment gateway integration\n"
                "- Kitchen display system sync\n"
                "- Predictive demand and prep planning\n"
                "- QR table sessions with live guest notifications\n"
                "- Multi-branch analytics\n\n"
                "Scalability & Expansion:\n"
                "The separated frontend/backend structure supports cleaner APIs, role-based modules and future database services.\n\n"
                "Conclusion:\n"
                "DINE AI turns scattered restaurant activity into one coordinated service rhythm."
            ),
        },
    ]

    for index, slide in enumerate(prs.slides):
        data = content[index]
        clear_template_content(slide)
        team_logo(slide)
        set_template_title(slide, data["title"])

        if index == 0:
            panel(slide, Inches(0.65), Inches(1.25), Inches(5.6), Inches(4.75))
            write_textbox(slide, data["body"], Inches(0.9), Inches(1.45), Inches(5.1), Inches(4.2), 18, INK)
            accent_panel(slide, "Core Promise", "Every service beat, in sync.", Inches(7.2), Inches(5.35), Inches(3.4), Inches(0.8), MINT)
        elif index in (2, 3, 4, 5):
            panel(slide, Inches(1.25), Inches(1.55), Inches(5.45), Inches(4.8))
            write_textbox(slide, data["body"], Inches(1.55), Inches(1.78), Inches(4.85), Inches(4.25), 14, INK)
        else:
            panel(slide, Inches(0.7), Inches(1.35), Inches(6.25), Inches(4.95))
            write_textbox(slide, data["body"], Inches(1.0), Inches(1.58), Inches(5.65), Inches(4.4), 14, INK)

        if index == 1:
            accent_panel(slide, "Service Gap", "Teams make decisions from different signals.", Inches(7.25), Inches(2.05), Inches(3.35), Inches(0.95), CORAL)
            accent_panel(slide, "Safety Gap", "Dietary context must stay visible through every handoff.", Inches(7.25), Inches(3.35), Inches(3.35), Inches(0.95), PURPLE)
        if index == 2:
            accent_panel(slide, "Unique Value", "One interface connects the guest promise to the operational work needed to keep it.", Inches(7.25), Inches(2.5), Inches(3.8), Inches(1.25), MINT)
        if index == 3:
            accent_panel(slide, "Clean Structure", "frontend/ for UI, backend/ for app.js, routes and controllers.", Inches(7.1), Inches(2.6), Inches(4.0), Inches(1.25), CORAL)
        if index == 4:
            accent_panel(slide, "Measurable Signals", "72 min table turn | 100% SafePlate checks | 8 stockouts avoided", Inches(7.0), Inches(2.65), Inches(4.1), Inches(1.25), MINT)
        if index == 5:
            accent_panel(slide, "Ready for Demo", "Frontend runs on localhost:3001 and backend exposes clean API routes.", Inches(7.1), Inches(2.65), Inches(4.0), Inches(1.25), CITRUS)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
