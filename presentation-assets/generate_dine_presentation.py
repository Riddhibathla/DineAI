from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(
    r"C:\Users\Samsung\Downloads\Vibeathon_6.0_Vibecoding_Hackathon_July_2026_Idea_Submission_Template.pptx"
)
OUTPUT = ROOT / "DINE_AI_Project_Pitch_Vibeathon_2026.pptx"

PURPLE = RGBColor(52, 29, 66)
CORAL = RGBColor(255, 111, 97)
MINT = RGBColor(114, 225, 190)
CITRUS = RGBColor(216, 250, 73)
INK = RGBColor(20, 21, 34)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text, size=22, color=INK, bold=False):
    shape.text = text
    frame = shape.text_frame
    frame.word_wrap = True
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def set_title(slide, text):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            set_text(shape, text, size=34, color=PURPLE, bold=True)
            return


def set_brand(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and "Team Name" in shape.text:
            set_text(shape, "DINE AI", size=18, color=PURPLE, bold=True)


def add_tag(slide, text, left, top, width, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = PURPLE if fill == CITRUS else WHITE


def add_card(slide, title, body, left, top, width, height, fill=WHITE, accent=CORAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(232, 229, 223)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent
    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.13), width - Inches(0.32), height - Inches(0.2))
    box.text = title + "\n" + body
    paragraphs = box.text_frame.paragraphs
    paragraphs[0].runs[0].font.bold = True
    paragraphs[0].runs[0].font.size = Pt(15)
    paragraphs[0].runs[0].font.color.rgb = PURPLE
    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.font.size = Pt(10.5)
            run.font.color.rgb = INK


def replace_body(slide, text, size=18):
    markers = (
        "TITLE PAGE",
        "Problem Overview",
        "Solution Overview",
        "Technologies Used",
        "Key Use Cases",
        "Future Enhancements",
    )
    for shape in slide.shapes:
        if shape.has_text_frame and any(marker in shape.text for marker in markers):
            set_text(shape, text, size=size, color=INK)
            return


def main():
    prs = Presentation(TEMPLATE)

    slide1, slide2, slide3, slide4, slide5, slide6 = prs.slides

    set_title(slide1, "DINE AI")
    replace_body(
        slide1,
        "TITLE PAGE\n\n"
        "Team Name: Byte Crunchers\n"
        "Team Leader Name: Riddhi Bathla\n"
        "College Name: Institute of Management Studies\n"
        "Year & Department: 3rd year & BCA\n\n"
        "Problem Statement / Project Title:\n"
        "DINE AI - Smart Restaurant Operations Platform\n\n"
        "A real-time operating layer for guests, service teams, kitchen, inventory, billing and dietary safety.",
        size=20,
    )
    add_tag(slide1, "LIVE OPS", Inches(8.2), Inches(5.8), Inches(1.25), CITRUS)
    add_tag(slide1, "SAFEPLATE", Inches(9.6), Inches(5.8), Inches(1.45), MINT)

    set_brand(slide2)
    set_title(slide2, "CURRENT PROBLEM")
    replace_body(
        slide2,
        "Problem Overview:\n"
        "Restaurants often run orders, waitlists, kitchen tickets, inventory and billing in separate tools or manual workflows.\n\n"
        "Who Is Affected:\n"
        "Guests, servers, kitchen staff, managers and owners.\n\n"
        "Pain Points:\n"
        "- Long waits and unclear table availability\n"
        "- Dietary requirements can get lost between guest, server and kitchen\n"
        "- Inventory shortages are noticed too late\n"
        "- Managers lack one live view of service health\n\n"
        "Why It Matters:\n"
        "A small communication miss can delay service, waste stock or create a safety risk.",
        size=17,
    )
    add_card(slide2, "Service gap", "Front-of-house and kitchen teams make decisions from different signals.", Inches(7.5), Inches(2.0), Inches(3.1), Inches(1.0), accent=CORAL)
    add_card(slide2, "Safety gap", "Dietary context needs to stay visible through every handoff.", Inches(7.5), Inches(3.25), Inches(3.1), Inches(1.0), accent=PURPLE)

    set_brand(slide3)
    set_title(slide3, "PROPOSED SOLUTION")
    replace_body(
        slide3,
        "Solution Overview:\n"
        "DINE AI brings the restaurant shift into one live workspace.\n\n"
        "How It Works:\n"
        "- Guests browse the menu, select preferences and send orders\n"
        "- Staff manage waitlist, tables and service tasks\n"
        "- Kitchen sees tickets with timing and safety context\n"
        "- Managers review operational signals and inventory pressure\n\n"
        "Key Features:\n"
        "- Guest menu and preference capture\n"
        "- Waitlist and service floor views\n"
        "- Kitchen line board\n"
        "- Pantry impact view\n"
        "- SafePlate dietary relay\n"
        "- Billing and analytics workspace\n\n"
        "Unique Value:\n"
        "One interface connects the guest promise to the operational work needed to keep it.",
        size=16,
    )
    add_card(slide3, "Guest", "Menu, wait time and preferences.", Inches(6.9), Inches(2.0), Inches(2.0), Inches(0.9), accent=MINT)
    add_card(slide3, "Kitchen", "Tickets with timing and safety state.", Inches(8.95), Inches(3.1), Inches(2.2), Inches(0.9), accent=CORAL)
    add_card(slide3, "Manager", "Signals, stock and shift health.", Inches(6.9), Inches(4.2), Inches(2.2), Inches(0.9), accent=PURPLE)

    set_brand(slide4)
    set_title(slide4, "TECHNICAL APPROACH")
    replace_body(
        slide4,
        "Technologies Used:\n"
        "- Frontend: Next.js, React, TypeScript, Tailwind CSS\n"
        "- Backend: Node.js app.js with route and controller separation\n"
        "- UI: Lucide icons, responsive CSS, demo data layer\n\n"
        "System Architecture:\n"
        "frontend/app -> frontend/components -> frontend/lib/demo-data\n"
        "backend/app.js -> backend/routes -> backend/controllers -> backend/data\n\n"
        "Methodology / Workflow:\n"
        "1. Route receives request\n"
        "2. Controller prepares domain response\n"
        "3. Frontend renders role-specific workspace\n"
        "4. Staff or guest acts from a focused operational view",
        size=15,
    )
    add_card(slide4, "frontend/", "Next app router pages and reusable workspace components.", Inches(7.0), Inches(2.0), Inches(3.4), Inches(0.9), accent=MINT)
    add_card(slide4, "backend/app.js", "HTTP entry point and API routing.", Inches(7.0), Inches(3.1), Inches(3.4), Inches(0.9), accent=CORAL)
    add_card(slide4, "routes + controllers", "Clean separation for API endpoints and business logic.", Inches(7.0), Inches(4.2), Inches(3.4), Inches(0.9), accent=PURPLE)

    set_brand(slide5)
    set_title(slide5, "USE CASES & IMPACT")
    replace_body(
        slide5,
        "Key Use Cases:\n"
        "- Guest ordering with dietary preferences\n"
        "- Staff table service and waitlist management\n"
        "- Kitchen order prioritization\n"
        "- Inventory visibility before stockouts\n"
        "- Manager shift analytics\n\n"
        "Target Users / Beneficiaries:\n"
        "Restaurant managers, servers, kitchen teams, guests and owners.\n\n"
        "Real-World Applications:\n"
        "Casual dining, cloud kitchens, cafes, hotel restaurants and multi-location operators.\n\n"
        "Expected Impact:\n"
        "- Faster service handoffs\n"
        "- Better allergen communication\n"
        "- Fewer stock surprises\n"
        "- Clearer manager decisions",
        size=15,
    )
    add_card(slide5, "72 min", "Target table turn insight.", Inches(7.1), Inches(2.05), Inches(1.5), Inches(0.9), accent=MINT)
    add_card(slide5, "100%", "SafePlate checks visible.", Inches(8.8), Inches(2.05), Inches(1.65), Inches(0.9), accent=PURPLE)
    add_card(slide5, "8", "Stockouts avoided signal.", Inches(7.95), Inches(3.25), Inches(1.65), Inches(0.9), accent=CORAL)

    set_brand(slide6)
    set_title(slide6, "FUTURE SCOPE & CONCLUSION")
    replace_body(
        slide6,
        "Future Enhancements:\n"
        "- Real POS integration\n"
        "- Kitchen display system sync\n"
        "- Predictive demand and prep planning\n"
        "- QR table sessions and live guest notifications\n"
        "- Multi-branch analytics\n\n"
        "Integration Opportunities:\n"
        "Payment gateways, inventory suppliers, reservation tools and loyalty platforms.\n\n"
        "Scalability & Expansion:\n"
        "The separated frontend/backend structure supports cleaner APIs, role-based modules and future database services.\n\n"
        "Conclusion:\n"
        "DINE AI turns scattered restaurant activity into one coordinated service rhythm.",
        size=15,
    )
    add_tag(slide6, "READY FOR DEMO", Inches(7.2), Inches(5.55), Inches(1.75), CITRUS)
    add_tag(slide6, "CLEAN ARCHITECTURE", Inches(9.15), Inches(5.55), Inches(2.1), MINT)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
