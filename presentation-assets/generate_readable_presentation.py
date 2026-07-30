from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "DINE_AI_Project_Pitch_Vibeathon_2026_READABLE.pptx"

PURPLE = RGBColor(45, 24, 61)
PLUM = RGBColor(84, 42, 103)
CORAL = RGBColor(230, 78, 82)
MINT = RGBColor(34, 166, 139)
CITRUS = RGBColor(210, 244, 73)
PAPER = RGBColor(255, 252, 246)
INK = RGBColor(24, 24, 32)
MUTED = RGBColor(86, 91, 105)
WHITE = RGBColor(255, 255, 255)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def textbox(slide, text, left, top, width, height, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, text, subtitle=None):
    textbox(slide, text, Inches(0.65), Inches(0.35), Inches(12.0), Inches(0.65), 32, PURPLE, True)
    if subtitle:
        textbox(slide, subtitle, Inches(0.72), Inches(1.03), Inches(11.5), Inches(0.35), 14, MUTED)


def pill(slide, text, left, top, width, fill=CITRUS, color=PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = color
    return shape


def card(slide, heading, bullets, left, top, width, height, accent=CORAL):
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    base.fill.solid()
    base.fill.fore_color.rgb = WHITE
    base.line.color.rgb = RGBColor(230, 224, 232)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.09), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent

    textbox(slide, heading, left + Inches(0.25), top + Inches(0.18), width - Inches(0.45), Inches(0.35), 18, PURPLE, True)
    body = "\n".join(f"- {item}" for item in bullets)
    textbox(slide, body, left + Inches(0.25), top + Inches(0.62), width - Inches(0.45), height - Inches(0.75), 14, INK)


def metric(slide, number, label, left, top, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.4), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    textbox(slide, number, left + Inches(0.15), top + Inches(0.12), Inches(2.1), Inches(0.42), 25, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, label, left + Inches(0.15), top + Inches(0.62), Inches(2.1), Inches(0.32), 11, WHITE, False, PP_ALIGN.CENTER)


def arrow(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.55), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CORAL
    shape.line.color.rgb = CORAL


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = blank_slide(prs)
    fill_background(slide, PURPLE)
    textbox(slide, "DINE AI", Inches(0.8), Inches(0.75), Inches(7.0), Inches(0.9), 48, WHITE, True)
    textbox(slide, "Smart Restaurant Operations Platform", Inches(0.85), Inches(1.65), Inches(7.5), Inches(0.4), 22, CITRUS, True)
    textbox(slide, "A live operating layer for guests, service teams, kitchen, inventory, billing and dietary safety.", Inches(0.85), Inches(2.35), Inches(6.2), Inches(1.0), 20, WHITE)
    card(slide, "Hackathon Details", [
        "Vibeathon 6.0 - Vibecoding Hackathon, July 2026",
        "Team: Byte Crunchers",
        "Team Leader: Riddhi Bathla",
        "College: Institute of Management Studies",
        "Year & Department: 3rd year & BCA",
    ], Inches(7.55), Inches(1.15), Inches(4.85), Inches(3.1), MINT)
    pill(slide, "LIVE OPS", Inches(0.85), Inches(5.65), Inches(1.45))
    pill(slide, "SAFEPLATE", Inches(2.55), Inches(5.65), Inches(1.65), MINT, WHITE)

    slide = blank_slide(prs)
    fill_background(slide, PAPER)
    title(slide, "CURRENT PROBLEM", "Restaurants lose time and safety context when operations are split across manual tools.")
    card(slide, "Operational Friction", [
        "Orders, waitlists, kitchen tickets and inventory run separately",
        "Managers do not get one live picture of the shift",
        "Staff decisions depend on scattered updates",
    ], Inches(0.7), Inches(1.65), Inches(5.8), Inches(2.05), CORAL)
    card(slide, "Guest Experience Gaps", [
        "Wait times feel unclear",
        "Dietary preferences can get missed during handoff",
        "Menu availability changes are not always visible fast enough",
    ], Inches(6.8), Inches(1.65), Inches(5.8), Inches(2.05), PLUM)
    card(slide, "Why It Matters", [
        "A small communication miss can delay service",
        "A missed allergy note can create real safety risk",
        "Late inventory awareness causes stockouts and unhappy guests",
    ], Inches(0.7), Inches(4.15), Inches(11.9), Inches(1.85), MINT)

    slide = blank_slide(prs)
    fill_background(slide, PAPER)
    title(slide, "PROPOSED SOLUTION", "DINE AI connects every role in one live restaurant workspace.")
    steps = [("Guest", "Menu, wait time and preferences"), ("Service", "Tables, tasks and handoffs"), ("Kitchen", "Tickets, timing and safety"), ("Manager", "Signals, stock and billing")]
    x = Inches(0.8)
    for index, (head, body) in enumerate(steps):
        card(slide, head, [body], x + Inches(index * 3.0), Inches(2.05), Inches(2.35), Inches(1.55), [MINT, CORAL, PLUM, MINT][index])
        if index < len(steps) - 1:
            arrow(slide, x + Inches(index * 3.0 + 2.43), Inches(2.65))
    card(slide, "Key Features", [
        "Guest menu and preference capture",
        "Waitlist and service floor views",
        "Kitchen line board",
        "Pantry impact view",
        "SafePlate dietary relay",
        "Billing and analytics workspace",
    ], Inches(1.15), Inches(4.4), Inches(10.9), Inches(1.8), CORAL)

    slide = blank_slide(prs)
    fill_background(slide, PAPER)
    title(slide, "TECHNICAL APPROACH", "Clean frontend/backend separation with route and controller layers.")
    card(slide, "Frontend", [
        "Next.js App Router",
        "React + TypeScript",
        "Tailwind CSS styling",
        "Reusable workspace components",
    ], Inches(0.8), Inches(1.65), Inches(3.8), Inches(2.35), MINT)
    card(slide, "Backend", [
        "Node.js app.js entry point",
        "Separated routes and controllers",
        "JSON endpoints for menu, operations and dashboard",
        "Demo data isolated under backend/data",
    ], Inches(4.85), Inches(1.65), Inches(3.8), Inches(2.35), CORAL)
    card(slide, "Workflow", [
        "Route receives request",
        "Controller prepares response",
        "Frontend renders role workspace",
        "User acts from a focused view",
    ], Inches(8.9), Inches(1.65), Inches(3.8), Inches(2.35), PLUM)
    textbox(slide, "frontend/app -> components -> demo data     |     backend/app.js -> routes -> controllers -> data", Inches(1.25), Inches(5.1), Inches(10.8), Inches(0.45), 18, PURPLE, True, PP_ALIGN.CENTER)

    slide = blank_slide(prs)
    fill_background(slide, PAPER)
    title(slide, "USE CASES & IMPACT", "Built for practical restaurant workflows during a live shift.")
    card(slide, "Use Cases", [
        "Guest ordering with dietary preferences",
        "Staff table service and waitlist management",
        "Kitchen prioritization",
        "Inventory visibility before stockouts",
        "Manager shift analytics",
    ], Inches(0.8), Inches(1.6), Inches(5.2), Inches(3.15), CORAL)
    card(slide, "Target Users", [
        "Restaurant managers",
        "Servers and hosts",
        "Kitchen staff",
        "Guests",
        "Owners and operators",
    ], Inches(6.35), Inches(1.6), Inches(5.2), Inches(3.15), MINT)
    metric(slide, "72 min", "table turn signal", Inches(1.4), Inches(5.35), PLUM)
    metric(slide, "100%", "SafePlate visibility", Inches(5.45), Inches(5.35), MINT)
    metric(slide, "8", "stockouts avoided", Inches(9.5), Inches(5.35), CORAL)

    slide = blank_slide(prs)
    fill_background(slide, PURPLE)
    textbox(slide, "FUTURE SCOPE & CONCLUSION", Inches(0.75), Inches(0.55), Inches(10.5), Inches(0.65), 32, WHITE, True)
    card(slide, "Future Enhancements", [
        "POS and payment gateway integrations",
        "Kitchen display system sync",
        "Predictive demand and prep planning",
        "QR table sessions with live notifications",
        "Multi-branch analytics",
    ], Inches(0.85), Inches(1.65), Inches(5.55), Inches(3.4), MINT)
    card(slide, "Conclusion", [
        "DINE AI turns scattered restaurant activity into one coordinated service rhythm",
        "The clean architecture makes future database, auth and integration work easier",
        "The project is ready for demo on localhost:3001",
    ], Inches(6.85), Inches(1.65), Inches(5.55), Inches(3.4), CORAL)
    pill(slide, "READY FOR DEMO", Inches(3.9), Inches(6.05), Inches(2.2))
    pill(slide, "CLEAN ARCHITECTURE", Inches(6.35), Inches(6.05), Inches(2.6), MINT, WHITE)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
